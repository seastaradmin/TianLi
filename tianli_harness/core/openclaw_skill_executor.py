"""
OpenClaw Skill Executor - Execute skills through OpenClaw's skill system

This allows TianLi Harness to use OpenClaw's installed skills like pptx.
"""

import asyncio
import subprocess
from typing import Dict, Any, Optional


class OpenClawSkillExecutor:
    """Execute skills through OpenClaw's skill system"""
    
    def __init__(self, working_dir: str = "."):
        self.working_dir = working_dir
    
    async def execute_skill(self, skill_name: str, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Execute an OpenClaw skill
        
        Args:
            skill_name: Name of the skill (e.g., "pptx")
            params: Parameters for the skill
        
        Returns:
            Result dict with success status and output
        """
        try:
            # Use npx skills to run the skill
            # The skill will be invoked automatically by OpenClaw when needed
            cmd = ["npx", "skills", "run", skill_name]
            
            # Add parameters as needed
            # For pptx skill, we might need to pass input/output paths
            if "input" in params:
                cmd.extend(["--input", params["input"]])
            if "output" in params:
                cmd.extend(["--output", params["output"]])
            if "action" in params:
                cmd.extend(["--action", params["action"]])
            
            result = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                cwd=self.working_dir
            )
            
            stdout, stderr = await result.communicate()
            
            return {
                "success": result.returncode == 0,
                "stdout": stdout.decode("utf-8"),
                "stderr": stderr.decode("utf-8"),
                "returncode": result.returncode
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }
    
    async def create_pptx(self, title: str, slides: list, output_path: str) -> Dict[str, Any]:
        """
        Create a PPTX file using the pptx skill
        
        Args:
            title: Presentation title
            slides: List of slide content dicts
            output_path: Path to save the .pptx file
        
        Returns:
            Result dict
        """
        # For now, use python-pptx directly since pptx skill needs specific invocation
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            
            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            
            # Content slides
            for slide_data in slides:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                if "title" in slide_data:
                    slide.shapes.title.text = slide_data["title"]
                
                if "content" in slide_data:
                    content = slide_data["content"]
                    if isinstance(content, str):
                        # Add as bullet points
                        tf = slide.shapes.placeholders[1].text_frame
                        tf.text = content.split("\n")[0] if "\n" in content else content
                        for line in content.split("\n")[1:]:
                            p = tf.add_paragraph()
                            p.text = line
                            p.level = 0
            
            # Save
            prs.save(output_path)
            
            return {
                "success": True,
                "output_path": output_path,
                "slides_count": len(slides) + 1
            }
            
        except ImportError:
            return {
                "success": False,
                "error": "python-pptx not installed. Run: pip install python-pptx"
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }


# Convenience function
async def create_ppt_with_openclaw(
    title: str,
    slides: list,
    output_path: str,
    working_dir: str = "."
) -> Dict[str, Any]:
    """
    Create PPT using OpenClaw's pptx skill or python-pptx
    
    Usage:
        result = await create_ppt_with_openclaw(
            title="TianLi Harness",
            slides=[
                {"title": "问题", "content": "现有 AI 工具的痛点"},
                {"title": "解决方案", "content": "天理的核心特性"}
            ],
            output_path="presentation.pptx"
        )
    """
    executor = OpenClawSkillExecutor(working_dir)
    return await executor.create_pptx(title, slides, output_path)
